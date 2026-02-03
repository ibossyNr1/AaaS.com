#!/usr/bin/env python3
"""
Pitch Deck Generator - Generate PPTX from YAML data

This script generates a professional pitch deck PPTX from a template and data file.

Usage:
    python3 generate_pptx.py --data company_data.yaml
    python3 generate_pptx.py --data company_data.yaml --template template.pptx
    python3 generate_pptx.py --data company_data.yaml --output my_deck.pptx

Requirements:
    pip install python-pptx pyyaml
"""

import argparse
import yaml
import sys
from pathlib import Path
from datetime import datetime

try:
    from pptx import Presentation
    from pptx.util import Inches, Pt
    from pptx.dml.color import RgbColor
    from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
    from pptx.enum.shapes import MSO_SHAPE
except ImportError:
    print("❌ python-pptx not installed. Run: pip install python-pptx pyyaml")
    sys.exit(1)


# Default colors (can be customized)
COLORS = {
    "primary": RgbColor(0x00, 0x66, 0xCC),      # Blue
    "secondary": RgbColor(0x00, 0x99, 0x66),    # Green
    "dark": RgbColor(0x33, 0x33, 0x33),         # Dark gray
    "light": RgbColor(0xF5, 0xF5, 0xF5),        # Light gray
    "white": RgbColor(0xFF, 0xFF, 0xFF),        # White
    "accent": RgbColor(0xFF, 0x66, 0x00),       # Orange
}


def load_data(data_path: Path) -> dict:
    """Load pitch data from YAML file."""
    with open(data_path, 'r', encoding='utf-8') as f:
        return yaml.safe_load(f)


def add_title_slide(prs: Presentation, data: dict):
    """Add title/cover slide."""
    slide_layout = prs.slide_layouts[6]  # Blank layout
    slide = prs.slides.add_slide(slide_layout)
    
    company = data.get('company', {})
    title_data = data.get('title', {})
    
    # Company name
    left, top, width, height = Inches(1), Inches(2.5), Inches(8), Inches(1.5)
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.text = company.get('name', 'Company Name')
    p.font.size = Pt(54)
    p.font.bold = True
    p.alignment = PP_ALIGN.CENTER
    
    # Tagline
    left, top, width, height = Inches(1), Inches(4), Inches(8), Inches(0.8)
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.text = company.get('tagline', '')
    p.font.size = Pt(24)
    p.font.italic = True
    p.alignment = PP_ALIGN.CENTER
    
    # Traction headline (if provided)
    headline = title_data.get('headline', '')
    if headline:
        left, top, width, height = Inches(1), Inches(4.8), Inches(8), Inches(0.5)
        txBox = slide.shapes.add_textbox(left, top, width, height)
        tf = txBox.text_frame
        p = tf.paragraphs[0]
        p.text = headline
        p.font.size = Pt(18)
        p.font.bold = True
        p.alignment = PP_ALIGN.CENTER


def add_problem_slide(prs: Presentation, data: dict):
    """Add problem statement slide."""
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)
    
    problem = data.get('problem', {})
    
    # Title
    add_slide_title(slide, problem.get('headline', 'The Problem'))
    
    # Pain point
    left, top, width, height = Inches(0.5), Inches(1.8), Inches(9), Inches(1)
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.text = problem.get('pain_point', '')
    p.font.size = Pt(20)
    p.font.bold = True
    
    # Details
    details = [
        f"👥 {problem.get('who_affected', '')}",
        f"❌ {problem.get('current_solutions', '')}",
        f"💰 {problem.get('cost_of_problem', '')}",
    ]
    
    top = Inches(2.8)
    for detail in details:
        if detail.split(' ', 1)[1]:  # Only add if there's content after emoji
            left, width, height = Inches(0.5), Inches(9), Inches(0.6)
            txBox = slide.shapes.add_textbox(left, top, width, height)
            tf = txBox.text_frame
            p = tf.paragraphs[0]
            p.text = detail
            p.font.size = Pt(18)
            top += Inches(0.7)


def add_solution_slide(prs: Presentation, data: dict):
    """Add solution slide."""
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)
    
    solution = data.get('solution', {})
    
    add_slide_title(slide, solution.get('headline', 'Our Solution'))
    
    # One-liner
    left, top, width, height = Inches(0.5), Inches(1.8), Inches(9), Inches(0.8)
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.text = solution.get('one_liner', '')
    p.font.size = Pt(22)
    p.font.bold = True
    
    # Key benefits
    benefits = solution.get('key_benefits', [])
    top = Inches(3)
    for benefit in benefits:
        left, width, height = Inches(0.5), Inches(9), Inches(0.5)
        txBox = slide.shapes.add_textbox(left, top, width, height)
        tf = txBox.text_frame
        p = tf.paragraphs[0]
        p.text = f"✓ {benefit}"
        p.font.size = Pt(18)
        top += Inches(0.6)


def add_market_slide(prs: Presentation, data: dict):
    """Add market opportunity slide with TAM/SAM/SOM."""
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)
    
    market = data.get('market', {})
    
    add_slide_title(slide, market.get('headline', 'Market Opportunity'))
    
    # TAM/SAM/SOM
    metrics = [
        ("TAM", market.get('tam', ''), market.get('tam_description', '')),
        ("SAM", market.get('sam', ''), market.get('sam_description', '')),
        ("SOM", market.get('som', ''), market.get('som_description', '')),
    ]
    
    left = Inches(0.5)
    for i, (label, value, desc) in enumerate(metrics):
        top = Inches(2 + i * 1.3)
        
        # Label
        txBox = slide.shapes.add_textbox(left, top, Inches(1), Inches(0.4))
        tf = txBox.text_frame
        p = tf.paragraphs[0]
        p.text = label
        p.font.size = Pt(16)
        p.font.bold = True
        
        # Value
        txBox = slide.shapes.add_textbox(Inches(1.5), top, Inches(2), Inches(0.4))
        tf = txBox.text_frame
        p = tf.paragraphs[0]
        p.text = value
        p.font.size = Pt(28)
        p.font.bold = True
        
        # Description
        txBox = slide.shapes.add_textbox(Inches(3.5), top, Inches(6), Inches(0.4))
        tf = txBox.text_frame
        p = tf.paragraphs[0]
        p.text = desc
        p.font.size = Pt(14)


def add_traction_slide(prs: Presentation, data: dict):
    """Add traction slide."""
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)
    
    traction = data.get('traction', {})
    
    add_slide_title(slide, traction.get('headline', 'Traction & Growth'))
    
    # Key metrics
    growth_metrics = traction.get('growth_metrics', [])
    left = Inches(0.5)
    for i, metric in enumerate(growth_metrics[:3]):
        top = Inches(2 + i * 1.2)
        
        txBox = slide.shapes.add_textbox(left, top, Inches(9), Inches(0.8))
        tf = txBox.text_frame
        p = tf.paragraphs[0]
        p.text = f"{metric.get('metric', '')}: {metric.get('current', '')} ({metric.get('growth', '')})"
        p.font.size = Pt(22)
        p.font.bold = True


def add_team_slide(prs: Presentation, data: dict):
    """Add team slide."""
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)
    
    team = data.get('team', {})
    
    add_slide_title(slide, team.get('headline', 'The Team'))
    
    founders = team.get('founders', [])
    for i, founder in enumerate(founders[:3]):
        left = Inches(0.5 + i * 3)
        top = Inches(2)
        
        # Name and title
        txBox = slide.shapes.add_textbox(left, top, Inches(3), Inches(0.5))
        tf = txBox.text_frame
        p = tf.paragraphs[0]
        p.text = founder.get('name', '')
        p.font.size = Pt(18)
        p.font.bold = True
        
        txBox = slide.shapes.add_textbox(left, Inches(2.5), Inches(3), Inches(0.4))
        tf = txBox.text_frame
        p = tf.paragraphs[0]
        p.text = founder.get('title', '')
        p.font.size = Pt(14)
        
        # Background
        txBox = slide.shapes.add_textbox(left, Inches(3), Inches(2.8), Inches(1.5))
        tf = txBox.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = founder.get('background', '')
        p.font.size = Pt(12)


def add_ask_slide(prs: Presentation, data: dict):
    """Add the ask slide."""
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)
    
    ask = data.get('the_ask', {})
    
    add_slide_title(slide, ask.get('headline', 'The Ask'))
    
    # Amount and round
    left, top, width, height = Inches(0.5), Inches(2), Inches(9), Inches(1)
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.text = f"{ask.get('amount', '')} {ask.get('round_type', '')} Round"
    p.font.size = Pt(36)
    p.font.bold = True
    p.alignment = PP_ALIGN.CENTER
    
    # Use of funds
    use_of_funds = ask.get('use_of_funds', [])
    top = Inches(3.5)
    for item in use_of_funds:
        left, width, height = Inches(0.5), Inches(9), Inches(0.5)
        txBox = slide.shapes.add_textbox(left, top, width, height)
        tf = txBox.text_frame
        p = tf.paragraphs[0]
        p.text = f"• {item.get('category', '')}: {item.get('percentage', '')}% - {item.get('details', '')}"
        p.font.size = Pt(16)
        top += Inches(0.5)
    
    # Milestones
    milestones = ask.get('milestones', [])
    if milestones:
        top += Inches(0.5)
        txBox = slide.shapes.add_textbox(Inches(0.5), top, Inches(9), Inches(0.4))
        tf = txBox.text_frame
        p = tf.paragraphs[0]
        p.text = "Milestones: " + " → ".join(milestones)
        p.font.size = Pt(14)
        p.font.bold = True


def add_slide_title(slide, title: str):
    """Add a title to a slide."""
    left, top, width, height = Inches(0.5), Inches(0.5), Inches(9), Inches(1)
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(32)
    p.font.bold = True


def generate_deck(data: dict, template_path: Path = None, output_path: Path = None):
    """Generate the complete pitch deck."""
    
    # Create presentation (from template if provided)
    if template_path and template_path.exists():
        prs = Presentation(template_path)
        print(f"📄 Using template: {template_path}")
    else:
        prs = Presentation()
        prs.slide_width = Inches(10)
        prs.slide_height = Inches(7.5)
        print("📄 Creating blank presentation")
    
    # Generate slides
    print("📊 Generating slides...")
    
    add_title_slide(prs, data)
    print("  ✓ Title slide")
    
    add_problem_slide(prs, data)
    print("  ✓ Problem slide")
    
    add_solution_slide(prs, data)
    print("  ✓ Solution slide")
    
    add_market_slide(prs, data)
    print("  ✓ Market slide")
    
    add_traction_slide(prs, data)
    print("  ✓ Traction slide")
    
    add_team_slide(prs, data)
    print("  ✓ Team slide")
    
    add_ask_slide(prs, data)
    print("  ✓ Ask slide")
    
    # Save
    if not output_path:
        company_name = data.get('company', {}).get('name', 'pitch_deck')
        safe_name = "".join(c for c in company_name if c.isalnum() or c in ' -_').strip()
        timestamp = datetime.now().strftime("%Y%m%d")
        output_path = Path(f"{safe_name}_pitch_deck_{timestamp}.pptx")
    
    prs.save(output_path)
    print(f"\n✅ Saved: {output_path}")
    print(f"   Slides: {len(prs.slides)}")
    
    return output_path


def main():
    parser = argparse.ArgumentParser(
        description="Generate pitch deck PPTX from YAML data"
    )
    parser.add_argument("--data", "-d", required=True, type=Path,
                        help="Path to YAML data file")
    parser.add_argument("--template", "-t", type=Path,
                        help="Path to template PPTX (optional)")
    parser.add_argument("--output", "-o", type=Path,
                        help="Output PPTX path (optional)")
    
    args = parser.parse_args()
    
    if not args.data.exists():
        print(f"❌ Data file not found: {args.data}")
        sys.exit(1)
    
    print(f"\n🚀 Pitch Deck Generator")
    print(f"=" * 40)
    
    data = load_data(args.data)
    print(f"📋 Loaded data: {args.data}")
    
    output = generate_deck(
        data, 
        template_path=args.template,
        output_path=args.output
    )
    
    print(f"\n💡 Next steps:")
    print(f"   1. Open {output} in PowerPoint/Keynote/Google Slides")
    print(f"   2. Add screenshots, logos, and charts")
    print(f"   3. Apply your brand styling")
    print(f"   4. Export to PDF for investor emails")


if __name__ == "__main__":
    main()
